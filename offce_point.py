from pptx import Presentation
prs = Presentation()
slide_layout = prs.slide_layouts[1]  
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Hello, OpenAI!"
subtitle.text = "Python-pptx is awesome!"
prs.save('test2.pptx')

# 作成したパワーポイントファイル内の文字数をカウント
total_characters = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            total_characters += len(shape.text)
print("合計文字数は:", total_characters)